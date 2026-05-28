"""News-Grasp アーキテクチャ概要 PPTX 生成スクリプト（最終版・D 案運用）。

実行: python build_architecture_pptx.py
出力: ./architecture.pptx（PowerPoint で開いていれば連番フォールバック）
"""
from __future__ import annotations

import os
import sys

NRI_SKILL = r"C:\Users\hidek\.claude\skills\nri-ppt-template\assets"
sys.path.insert(0, NRI_SKILL)

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.oxml.ns import qn

import layouts  # type: ignore
from theme import (  # type: ignore
    NRI_THEME,
    set_shape_fill,
    set_shape_line,
    set_ja_font,
)

ORG = "News-Grasp"
# 出力先は Obsidian ボルト内の本物 repo の docs/
REPO_DOCS = r"C:\Users\hidek\Obsidian\New's Grasp\News-Grasp\docs"
os.makedirs(REPO_DOCS, exist_ok=True)
_BASE_PATH = os.path.join(REPO_DOCS, "architecture.pptx")
OUT_PATH = _BASE_PATH

CATEGORIES = [
    ("FX",            "Foreign Exchange",        (0xB8, 0x86, 0x0B), "¥"),
    ("AI",            "Artificial Intelligence", (0x2D, 0x5B, 0xB8), "◆"),
    ("IT-Consulting", "IT & Consulting",         (0x2E, 0x6B, 0x52), "▲"),
    ("Economy",       "Economy",                 (0x8E, 0x2A, 0x19), "■"),
    ("Game",          "Gaming",                  (0x5E, 0x3D, 0x8C), "●"),
]


# ---------------------------------------------------------------------------
# カスタム描画ヘルパー
# ---------------------------------------------------------------------------

def add_block(
    slide, text, *, x, y, w, h, fill, color,
    font_size=13, bold=True, rounded=True, line_color=None,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    set_shape_fill(shape, fill)
    set_shape_line(shape, line_color or NRI_THEME["ir_navy"], width_pt=0.75)
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.04)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    lines = text.split("\n")
    tf.text = lines[0]
    for ln in lines[1:]:
        tf.add_paragraph().text = ln
    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            set_ja_font(run, bold=bold, size_pt=font_size, color=color)
    return shape


def add_arrow(slide, *, x1, y1, x2, y2, color=None, width_pt=1.5, dash=False):
    if color is None:
        color = NRI_THEME["ir_navy"]
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    ln = line.line._get_or_add_ln()
    for tag in ("a:prstDash", "a:headEnd", "a:tailEnd"):
        existing = ln.find(qn(tag))
        if existing is not None:
            ln.remove(existing)
    if dash:
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))
    ln.append(ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"}))
    return line


def add_label(slide, text, *, x, y, w, h=0.3, font_size=11, color=None, align=PP_ALIGN.CENTER, bold=True):
    if color is None:
        color = NRI_THEME["ir_text_muted"]
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.text = text
    for p in tf.paragraphs:
        p.alignment = align
        for run in p.runs:
            set_ja_font(run, bold=bold, size_pt=font_size, color=color)
    return tb


# ---------------------------------------------------------------------------
# プレゼン本体
# ---------------------------------------------------------------------------

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

layouts.setup_slide_master(
    prs,
    org=ORG,
    disclaimer="本資料は News-Grasp プロジェクトの内部設計メモです。社外配布は想定していません。",
    show_page_number=True,
)

def _prepare_cover_image() -> str | None:
    """高解像度 cover-ai.png を JPEG に圧縮し、PPT に埋め込みやすいサイズにする。"""
    src = os.path.join(os.path.dirname(__file__), "cover-ai.png")
    if not os.path.exists(src):
        return None
    try:
        from PIL import Image
    except ImportError:
        return src
    out = os.path.join(os.path.dirname(__file__), "cover-ai-compressed.jpg")
    img = Image.open(src)
    img.thumbnail((1000, 1000))  # 縦横どちらか 1000px に収まる
    img.convert("RGB").save(out, "JPEG", quality=82, optimize=True)
    return out


COVER_IMAGE = _prepare_cover_image()


# ---- 1. 表紙 ----
layouts.add_cover_slide(
    prs,
    title="News-Grasp\n仕様 & アーキテクチャ",
    org=ORG,
    presenter_title="毎日 Web 情報収集 + 関連付け解説 Agent",
    presenter="",
    date="2026年4月28日 / 最終運用形態",
    tagline="Local Claude Code × GitHub × GAS Webhook · D-plan",
    cover_image=COVER_IMAGE,
)


# ---- 2. 目次 ----
layouts.add_agenda_slide(
    prs,
    items=[
        "全体アーキテクチャ（D 案）",
        "処理フロー（8 ステップ）",
        "デザインシステム & 5 カテゴリ",
        "Obsidian タグ仕様（階層タグ + entity）",
        "曜日マトリクス & サムネ運用",
        "採用方針の変遷（② → D 案）",
        "コスト・運用・テスト機構",
        "制約・リスク",
    ],
    page_num=1,
    org=ORG,
)


# ---- 3. 全体アーキテクチャ（D 案）----
slide3 = layouts.add_chart_slide(
    prs, category="アーキテクチャ", title="全体アーキテクチャ（ローカル Claude Code 主導）",
    page_num=2, org=ORG,
)

# Local PC 境界
boundary_pc = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.6), Inches(5.7), Inches(4.6),
)
set_shape_fill(boundary_pc, None)
set_shape_line(boundary_pc, NRI_THEME["ir_text_muted"], width_pt=0.75, dash=True)
add_label(slide3, "Local PC (Windows 11 / Max sub.)",
          x=0.5, y=1.45, w=4.0, h=0.25, font_size=10,
          color=NRI_THEME["ir_text_muted"], align=PP_ALIGN.LEFT)

# タスクスケジューラ
add_block(slide3, "Windows タスク\n「News-Grasp Runner」\n毎日 06:00 JST",
          x=0.7, y=1.85, w=2.4, h=1.0,
          fill=NRI_THEME["ir_blue_light"], color=NRI_THEME["ir_navy"], font_size=11)
# runner.bat
add_block(slide3, "news-grasp-runner.bat",
          x=3.4, y=1.85, w=2.4, h=0.5,
          fill=NRI_THEME["white"], color=NRI_THEME["ir_navy"], font_size=10)
add_label(slide3, "git pull → claude --print",
          x=3.4, y=2.4, w=2.4, h=0.3, font_size=9, color=NRI_THEME["ir_text_muted"])
# Claude Code
add_block(slide3, "claude.exe --print\nSonnet 4.6\n--tools default\n--dangerously-skip-permissions",
          x=0.7, y=3.1, w=5.1, h=1.3,
          fill=NRI_THEME["ir_navy"], color=NRI_THEME["white"], font_size=12)
# Obsidian
add_block(slide3, "Obsidian Vault\n  News's Grasp/News-Grasp/\n  digest/{Genre}/...",
          x=0.7, y=4.7, w=5.1, h=1.3,
          fill=NRI_THEME["ir_blue_light"], color=NRI_THEME["ir_navy"], font_size=11)

# クラウド境界
boundary_cloud = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.5), Inches(1.6), Inches(6.4), Inches(4.6),
)
set_shape_fill(boundary_cloud, None)
set_shape_line(boundary_cloud, NRI_THEME["ir_text_muted"], width_pt=0.75, dash=True)
add_label(slide3, "Cloud Services",
          x=6.6, y=1.45, w=2.0, h=0.25, font_size=10,
          color=NRI_THEME["ir_text_muted"], align=PP_ALIGN.LEFT)

# GitHub
add_block(slide3, "GitHub\nHIDEPON-UMG/News-Grasp\n(private repo)",
          x=6.8, y=1.85, w=2.7, h=1.3,
          fill=NRI_THEME["ir_blue_accent"], color=NRI_THEME["white"], font_size=12)
add_label(slide3, "digest/  data/  prompts/  assets/",
          x=6.8, y=3.18, w=2.7, h=0.3, font_size=9, color=NRI_THEME["ir_text_muted"])

# GAS Webhook
add_block(slide3, "GAS Web App\nnews-grasp-mailer\n(hidepontrainer@gmail)",
          x=10.0, y=1.85, w=2.7, h=1.3,
          fill=NRI_THEME["ir_navy_deep"], color=NRI_THEME["white"], font_size=12)
add_label(slide3, "client + 宛先ホワイトリスト",
          x=10.0, y=3.18, w=2.7, h=0.3, font_size=9, color=NRI_THEME["ir_text_muted"])

# Gmail（×2）
add_block(slide3, "Gmail\nhideki.kusunoki@gmail.com\nh2-hiramatsu@nri.co.jp",
          x=10.0, y=3.7, w=2.7, h=1.3,
          fill=NRI_THEME["ir_blue_light"], color=NRI_THEME["ir_navy"], font_size=11)

# 矢印
# Runner → Claude
add_arrow(slide3, x1=3.25, y1=2.85, x2=3.25, y2=3.1, color=NRI_THEME["ir_navy"])
# Claude ↔ GitHub: clone/pull/push
add_arrow(slide3, x1=5.85, y1=3.4, x2=6.8, y2=2.3, color=NRI_THEME["ir_navy"])
add_arrow(slide3, x1=6.8, y1=2.7, x2=5.85, y2=3.7, color=NRI_THEME["ir_navy"])
add_label(slide3, "git pull / push",
          x=5.6, y=2.55, w=1.6, h=0.25, font_size=9)
# Claude → GAS
add_arrow(slide3, x1=5.85, y1=3.55, x2=10.0, y2=2.3, color=NRI_THEME["ir_navy"])
add_label(slide3, "Webhook POST\n(client=news-grasp-routine)",
          x=7.5, y=2.55, w=2.0, h=0.4, font_size=9, color=NRI_THEME["ir_text_muted"])
# GAS → Gmail
add_arrow(slide3, x1=11.35, y1=3.18, x2=11.35, y2=3.7, color=NRI_THEME["ir_navy"])
add_label(slide3, "GmailApp.sendEmail",
          x=12.7, y=3.35, w=0.3, h=0.25, font_size=9)
# GitHub → Obsidian (runner.bat 内 git pull)
add_arrow(slide3, x1=6.8, y1=3.1, x2=5.85, y2=4.7,
          color=NRI_THEME["ir_navy_deep"], dash=True)
add_label(slide3, "Runner 内 git pull で同期",
          x=4.0, y=4.05, w=2.5, h=0.25, font_size=9, color=NRI_THEME["ir_text_muted"])


# ---- 4. 処理フロー（8 ステップ）----
slide4 = layouts.add_chart_slide(
    prs, category="フロー", title="処理フロー — 8 ステップ",
    page_num=3, org=ORG,
)
steps = [
    ("①", "曜日判定",  "対象カテゴリ\n(4 or 5)"),
    ("②", "状態 Read", "watchlist /\narticles.jsonl"),
    ("③", "WebSearch", "5 件 / カテゴリ\nスコア降順"),
    ("④", "OGP 取得",  "WebFetch x N\n失敗→NG画像"),
    ("⑤", "関連照合",  "過去90日 ×\n5 軸マッチ"),
    ("⑥", "digest 生成", "MD + JSONL\n(Summary 含)"),
    ("⑦", "commit/push", "per-category\nfolders"),
    ("⑧", "メール送信",  "GAS Webhook\n→ Gmail x 2"),
]
n = len(steps)
total_w = 12.4
gap = 0.12
box_w = (total_w - gap * (n - 1)) / n
box_h = 1.8
top = 2.6

for i, (idx, label, detail) in enumerate(steps):
    x = 0.45 + i * (box_w + gap)
    badge = slide4.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x + box_w/2 - 0.3), Inches(top - 0.6),
        Inches(0.6), Inches(0.6),
    )
    set_shape_fill(badge, NRI_THEME["ir_navy"])
    set_shape_line(badge, None)
    tf = badge.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.text = idx
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_ja_font(p.runs[0], bold=True, size_pt=18, color=NRI_THEME["white"])

    add_block(slide4, label, x=x, y=top, w=box_w, h=box_h * 0.55,
              fill=NRI_THEME["ir_blue_light"], color=NRI_THEME["ir_navy"], font_size=12)
    add_block(slide4, detail, x=x, y=top + box_h*0.6, w=box_w, h=box_h * 0.4,
              fill=NRI_THEME["white"], color=NRI_THEME["ir_text_muted"],
              font_size=8, bold=False)
    if i < n - 1:
        ax = x + box_w + 0.005
        ay = top + box_h * 0.275
        add_arrow(slide4, x1=ax, y1=ay, x2=ax + gap - 0.01, y2=ay,
                  color=NRI_THEME["ir_navy"], width_pt=2.0)

add_label(slide4,
    "失敗時は最大 3 回リトライ（30s → 60s → 120s）。最終失敗時は data/_status.md に追記し、Webhook で件名「[News-Grasp 失敗] YYYY-MM-DD」のメールを送信。",
    x=0.6, y=4.9, w=12.2, h=0.4, font_size=11, color=NRI_THEME["ir_text"], align=PP_ALIGN.LEFT)


# ---- 5. デザインシステム ----
layouts.add_data_table_slide(
    prs, category="デザイン", title="5 カテゴリ × アクセントカラー × 強調記法",
    headers=["ID", "カテゴリ", "英名", "アクセント色", "グリフ"],
    rows=[
        ["fx",      "為替",          "Foreign Exchange",        "#B8860B（琥珀）",      "¥"],
        ["ai",      "AI",            "Artificial Intelligence", "#2D5BB8（電子青）",    "◆"],
        ["it",      "IT-Consulting", "IT & Consulting",         "#2E6B52（苔緑）",      "▲"],
        ["economy", "経済",          "Economy",                  "#8E2A19（深紅）",      "■"],
        ["game",    "ゲーム",        "Gaming",                   "#5E3D8C（洋紫）",      "●"],
    ],
    commentary=[
        "強調記法 [[キーワード]] = 太字 + アクセント色背景（記事あたり 2〜4 箇所）",
        "強調記法 __重要文__ = 下線 + 太字（段落あたり 1〜2 箇所）",
        "本文 Noto Serif JP（明朝） / メタ JetBrains Mono / 欧文 Inter",
    ],
    page_num=4, org=ORG,
)


# ---- 6. Obsidian タグ仕様 ----
layouts.add_data_table_slide(
    prs, category="タグ仕様", title="Obsidian 階層タグ — 5 entity + 4 補助 + score",
    headers=["プレフィクス", "種別", "値の規則", "例"],
    rows=[
        ["（なし）",   "共通固定 4 件",       "英字、号番号入り",          "daily / newsletter / news-grasp / issue-{号}"],
        ["cat/",      "カテゴリ id",          "5 種固定",                 "cat/fx / cat/ai / cat/it / cat/economy / cat/game"],
        ["co/",       "企業／組織",            "日本語優先（英字固有名詞は原文）", "co/トヨタ / co/NTTデータ / co/OpenAI / co/NVIDIA"],
        ["country/",  "国",                   "日本語国名（EU は EU）",     "country/日本 / country/米国 / country/EU"],
        ["svc/",      "サービス／製品",         "原文。スペース→ハイフン、ピリオド→アンダースコア", "svc/Claude / svc/Switch-2 / svc/GPT-5_5"],
        ["person/",   "人名",                 "日本語フルネーム、中点 ・ OK",  "person/植田和男 / person/ジェローム・パウエル"],
        ["ticker/",   "株式・通貨",            "大文字。スラッシュは削除",     "ticker/USDJPY / ticker/NVDA / ticker/7974"],
        ["topic/",    "テーマ（1〜3 個）",      "日本語推奨、国際略号 OK",     "topic/利下げ / topic/FOMC / topic/規制"],
        ["industry/", "業界（0〜2 個）",        "日本語",                   "industry/半導体 / industry/IT-コンサル"],
        ["event/",    "イベント種別（0〜2 個）", "日本語",                   "event/決算 / event/製品発表 / event/政策会合"],
        ["score/",    "重要度",               "高 (≥85) / 中 (65-84) / 低 (<65)", "score/高 / score/中 / score/低"],
    ],
    commentary=[
        "タグ値で使えない文字：半角スペース → '-'、スラッシュ '/' → 削除、ピリオド '.' → '_'。中点 '・' と長音 'ー' は OK",
        "Runner は記事要約と同じターンで entities/topics/industries/events を JSON 出力 → tags[] に階層化",
        "frontmatter は cat/co/country/person のみ集約。記事カード行は 4〜7 個に絞る（圧縮版）。詳細: prompts/obsidian-tagging-spec.md",
    ],
    page_num=5, org=ORG,
)


# ---- 7. 曜日マトリクス & サムネ運用 ----
layouts.add_data_table_slide(
    prs, category="運用", title="曜日 × カテゴリマトリクス（毎朝 06:00 JST 起動）",
    headers=["曜日", "FX", "AI", "IT", "Economy", "Game", "計"],
    rows=[
        ["月", "●", "●", "●", "●", "—", "4"],
        ["火", "●", "●", "●", "●", "●", "5"],
        ["水", "●", "●", "●", "●", "—", "4"],
        ["木", "●", "●", "●", "●", "●", "5"],
        ["金", "●", "●", "●", "●", "—", "4"],
        ["土", "●", "●", "●", "—", "●", "4"],
        ["日", "●", "●", "●", "—", "●", "4"],
    ],
    commentary=[
        "FX は独立カテゴリで毎日掲載。Economy は平日のみ、Game は火木土日のみ",
        "サムネ：FEATURED 568×200 = カテゴリ別キービジュアル / サイド 140×90 = カテゴリ別共通系",
        "OGP 画像が取得できた記事はそれを優先表示、取得失敗時のみ NG プレースホルダ",
    ],
    page_num=6, org=ORG,
)


# ---- 8. 採用方針の変遷 ----
layouts.add_points_slide(
    prs, period="採用方針\nの変遷",
    sections=[
        {"title": "経緯", "bullets": [
            "計画時：4 案比較で ② Anthropic Routine を最有力に決定（Max 内・PC オフ非依存・実装最短）",
            "実装後：Routine の「クラウドコンテナをセットアップ中」が 1 時間以上ハングする事象が連続発生",
            "原因：新規環境のプロビジョニング不安定 + private repo の OAuth 連携が未対応",
            "翌朝の自動発火も同様に動かず、Routine 機能の実用適合性が運用に耐えないと判断",
        ]},
        {"title": "D 案（採用版）への移行と利点", "bullets": [
            "Windows タスクスケジューラ → claude.exe --print --tools default で本番運用",
            "Max サブスク内で完結（追加 API 課金 $0、5h 枠 15〜25%/回）",
            "GitHub repo・GAS Webhook・Obsidian ボルト・watchlist は ② から流用、無駄なし",
            "PC オン時のみ動くが、毎朝起動の習慣があれば実用上問題なし",
        ]},
    ],
    page_num=7, org=ORG,
)


# ---- 9. コスト ----
layouts.add_data_table_slide(
    prs, category="コスト", title="月額コスト内訳（追加課金ゼロ）",
    headers=["項目", "課金経路", "月額", "備考"],
    rows=[
        ["Claude Sonnet 4.6 実行", "Max サブスク内", "$0", "5h 枠の 15〜25%/回"],
        ["GitHub プライベート repo", "個人プラン無料枠", "$0", "Contents R/W、PAT 不要" ],
        ["GAS Webhook", "Workspace 無料枠", "$0", "Gmail 送信は 100 通/日まで"],
        ["メール配信", "Gmail (GAS 経由)", "$0", "送信元 hidepontrainer@gmail.com"],
        ["合計", "—", "$0", "Max 既存契約のみで完結"],
    ],
    current_col_index=2,
    commentary=[
        "5h 枠は朝 06:00 実行のためユーザーの作業時間と干渉しない",
        "5 カテゴリ × 5 件 + 5 セクション考察 + サムネ OGP 取得で約 50-70k トークン入力 / 6-9k トークン出力",
    ],
    page_num=8, org=ORG,
)


# ---- 10. ファイル構造 & テスト機構 ----
layouts.add_points_slide(
    prs, period="ファイル\n構造",
    sections=[
        {"title": "リポジトリ構成", "bullets": [
            "digest/{FX,AI,IT-Consulting,Economy,Game}/ — カテゴリ別フォルダで管理",
            "digest/Summary/ — 日次サマリー（テーマ考察ハブ）",
            "data/watchlist.md — ★ ユーザー編集可、トラッキング対象一覧",
            "data/articles.jsonl — 過去 90 日の記事メタ（自動ローテート、entities/topics/tags 拡張済み）",
            "prompts/routine-system.md / obsidian-tagging-spec.md / email-template.html / obsidian-template.md",
            "assets/ng-thumb-{cat}.jpg + ng-thumb-common-{cat}.jpg（v2 Claude Design 提供 10 枚）",
        ]},
        {"title": "単体テスト機構", "bullets": [
            "tests/render_email.py — Claude を呼ばずに HTML テンプレートだけ確認",
            "A) python tests/render_email.py — preview.html 生成（$0 / 1-2 秒）",
            "C) python tests/render_email.py --send — Webhook 経由で実送信（$0 / 5-10 秒）",
            "tests/mock_data.py に 5 カテゴリ × 5 件のサンプル（強調記法 [[]] / __ 含む）",
        ]},
    ],
    page_num=9, org=ORG,
)


# ---- 11. 制約・リスク ----
layouts.add_points_slide(
    prs, period="制約・\nリスク",
    sections=[
        {"title": "情報源・配信の制約", "bullets": [
            "NewsPicks 有料記事は認証ゲートのため公開部分のみ参照",
            "NRI 宛メール（h2-hiramatsu@nri.co.jp）は外部メールフィルタで弾かれる可能性",
            "→ 不達時は hideki.kusunoki@gmail.com で受信し、必要時に手動転送",
        ]},
        {"title": "運用上の依存・恒久ルール", "bullets": [
            "Max サブスクの 5h 枠を朝 6 時に消費（日中の作業と時間帯競合させない）",
            "PC オン時に Runner が動く前提（土日 PC オフだとその日はスキップされる）",
            "Windows .bat は CRLF 必須・ASCII のみ・goto 構造（feedback memory に恒久ルール記録済み）",
            "メール HTML 内画像は base64 inline 必須（private repo の raw URL は受信側からアクセス不可）",
        ]},
    ],
    page_num=10, org=ORG,
)


# ---- 12. 閉じ ----
layouts.add_closing_slide(
    prs,
    tagline="News-Grasp",
    subtitle="Five Lenses on Today · D-plan in production · 2026-04-28",
    page_num=None,
    org=ORG,
)


# ---- 書き込み（PowerPoint で開かれていれば連番フォールバック）----
out_path = OUT_PATH
for n in range(2, 10):
    try:
        prs.save(out_path)
        break
    except PermissionError:
        out_path = _BASE_PATH.replace(".pptx", f".v{n}.pptx")
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
